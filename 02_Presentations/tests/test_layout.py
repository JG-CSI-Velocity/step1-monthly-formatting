"""Tests for style.layout -- slide zones and image fitting."""

from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from style.layout import (
    CONTENT_ZONE,
    FOOTER_ZONE,
    KPI_ROW_H,
    SAFE_BOTTOM,
    SAFE_LEFT,
    SAFE_RIGHT,
    SAFE_TOP,
    SLIDE_H,
    SLIDE_W,
    TITLE_ZONE,
    Zone,
    fit_image,
    is_inside_zone,
)


def test_slide_is_13_33x7_5_widescreen():
    assert SLIDE_W == Inches(13.33)
    assert SLIDE_H == Inches(7.5)


def test_safe_margins():
    assert SAFE_LEFT == Inches(0.5)
    assert SAFE_RIGHT == Inches(0.5)
    assert SAFE_TOP == Inches(0.4)
    assert SAFE_BOTTOM == Inches(0.4)


def test_title_zone_spans_full_width_under_top_margin():
    assert TITLE_ZONE.top == Inches(0.4)
    assert TITLE_ZONE.height == Inches(1.0)
    assert TITLE_ZONE.full_width is True


def test_footer_zone_is_at_bottom():
    assert FOOTER_ZONE.height == Inches(0.4)


def test_content_zone_sits_between_title_and_footer():
    assert CONTENT_ZONE.top == Inches(1.4)
    assert CONTENT_ZONE.bottom_offset == Inches(0.8)


def test_kpi_row_height():
    assert KPI_ROW_H == Inches(1.2)


def test_zone_is_frozen_dataclass():
    z = Zone(top=Inches(1), height=Inches(1))
    with pytest.raises(Exception):
        z.top = Inches(2)


def _slide_with_shape(top_in: float, height_in: float):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_textbox(
        Inches(0.5), Inches(top_in), Inches(12), Inches(height_in)
    )
    return prs, shape


def test_shape_inside_title_zone():
    _, shape = _slide_with_shape(top_in=0.5, height_in=0.8)
    assert is_inside_zone(shape, TITLE_ZONE, slide_height=SLIDE_H) is True


def test_shape_outside_title_zone_below():
    _, shape = _slide_with_shape(top_in=2.0, height_in=0.8)
    assert is_inside_zone(shape, TITLE_ZONE, slide_height=SLIDE_H) is False


def test_shape_inside_footer_zone():
    _, shape = _slide_with_shape(top_in=7.15, height_in=0.3)
    assert is_inside_zone(shape, FOOTER_ZONE, slide_height=SLIDE_H) is True


def test_fit_image_preserves_aspect_ratio(tmp_path):
    img_path = tmp_path / "wide.png"
    Image.new("RGB", (1000, 500), (0xE7, 0x4C, 0x3C)).save(img_path)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    zone = Zone(top=Inches(1.4), bottom_offset=Inches(0.8), full_width=True,
                left=Inches(0.5), width=Inches(12.33))

    pic = fit_image(slide, img_path, zone, slide_height=Inches(7.5))

    assert abs(pic.width - Inches(12.33)) < Inches(0.01)
    assert abs(pic.height - Inches(6.165)) < Inches(0.02)
    assert pic.top == Inches(1.4)
