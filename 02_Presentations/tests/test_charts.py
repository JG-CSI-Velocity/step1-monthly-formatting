"""Tests for style.charts -- audit of embedded chart images."""

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from style.charts import audit_chart_image


def _slide_with_picture(tmp_path: Path, px_w: int, px_h: int, width_in: float):
    img_path = tmp_path / "chart.png"
    Image.new("RGB", (px_w, px_h), (0x0D, 0x94, 0x88)).save(img_path)
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pic = slide.shapes.add_picture(
        str(img_path),
        Inches(0.5),
        Inches(1.4),
        width=Inches(width_in),
        height=Inches(width_in * px_h / px_w),
    )
    return pic


def test_audit_meets_dpi_floor_when_resolution_is_high(tmp_path):
    # 1800 px rendered at 12" wide = 150 DPI
    pic = _slide_with_picture(tmp_path, px_w=1800, px_h=900, width_in=12.0)
    audit = audit_chart_image(pic)
    assert audit.dpi_estimate >= 150
    assert audit.meets_dpi_floor is True


def test_audit_below_dpi_floor_for_low_resolution(tmp_path):
    # 600 px at 12" wide = 50 DPI
    pic = _slide_with_picture(tmp_path, px_w=600, px_h=300, width_in=12.0)
    audit = audit_chart_image(pic)
    assert audit.dpi_estimate < 150
    assert audit.meets_dpi_floor is False
