"""Tests for style.typography -- Montserrat type scale per SLIDE_MAPPING.md."""

from pptx import Presentation
from pptx.util import Inches, Pt

from style.typography import (
    TextStyle,
    apply,
    axis_label,
    chart_subtitle,
    chart_title,
    data_label,
    footnote,
    kpi_hero,
    kpi_label,
    slide_title,
    subtitle,
)


def test_slide_title_is_20pt_montserrat_extrabold():
    s = slide_title()
    assert s == TextStyle(font_name="Montserrat", size_pt=20, weight="ExtraBold")


def test_subtitle_is_14pt_montserrat_regular():
    assert subtitle() == TextStyle("Montserrat", 14, "Regular")


def test_chart_title_is_18pt_montserrat_bold():
    assert chart_title() == TextStyle("Montserrat", 18, "Bold")


def test_chart_subtitle_is_12pt_montserrat_regular():
    assert chart_subtitle() == TextStyle("Montserrat", 12, "Regular")


def test_data_label_is_11pt_montserrat_regular():
    assert data_label() == TextStyle("Montserrat", 11, "Regular")


def test_axis_label_is_10pt_montserrat_regular():
    assert axis_label() == TextStyle("Montserrat", 10, "Regular")


def test_footnote_is_9pt_montserrat_regular():
    assert footnote() == TextStyle("Montserrat", 9, "Regular")


def test_kpi_hero_is_36pt_montserrat_extrabold():
    assert kpi_hero() == TextStyle("Montserrat", 36, "ExtraBold")


def test_kpi_label_is_11pt_montserrat_regular():
    assert kpi_label() == TextStyle("Montserrat", 11, "Regular")


def _one_slide_with_textbox(text: str):
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    tf = box.text_frame
    tf.text = text
    return prs, tf


def test_apply_sets_font_name_and_size_on_every_run():
    _, tf = _one_slide_with_textbox("Hello")
    apply(tf, slide_title())
    for para in tf.paragraphs:
        for run in para.runs:
            assert run.font.name == "Montserrat"
            assert run.font.size == Pt(20)
            assert run.font.bold is True  # ExtraBold -> bold


def test_apply_regular_weight_clears_bold():
    _, tf = _one_slide_with_textbox("Hello")
    apply(tf, subtitle())
    for para in tf.paragraphs:
        for run in para.runs:
            assert run.font.bold is False


def test_apply_handles_multiple_paragraphs_and_runs():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(2))
    tf = box.text_frame
    tf.text = "First"
    tf.add_paragraph().text = "Second"
    apply(tf, footnote())
    sizes = [run.font.size for p in tf.paragraphs for run in p.runs]
    assert all(s == Pt(9) for s in sizes)
