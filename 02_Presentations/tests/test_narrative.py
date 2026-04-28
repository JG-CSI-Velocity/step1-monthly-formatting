"""Tests for style.narrative -- 0-3 scorer on consultative/performance/focal axes."""

from pathlib import Path

from pptx import Presentation

from style.narrative import NarrativeScore, score_slide

FIXTURES = Path(__file__).parent / "fixtures"


def test_pristine_content_slide_scores_3_3_3():
    prs = Presentation(str(FIXTURES / "pristine.pptx"))
    slide = prs.slides[1]  # slide 1 is title, slides 2+ are content
    score = score_slide(slide, slide_height=prs.slide_height)
    assert score.consultative == 3
    assert score.performance == 3
    assert score.focal == 3


def test_fragment_headline_scores_low_consultative():
    prs = Presentation(str(FIXTURES / "moderately_broken.pptx"))
    slide = prs.slides[1]  # fragment headline slide
    score = score_slide(slide, slide_height=prs.slide_height)
    assert score.consultative <= 1


def test_missing_annotation_drops_performance_to_2():
    prs = Presentation(str(FIXTURES / "moderately_broken.pptx"))
    slide = prs.slides[2]  # note-only, no annotation
    score = score_slide(slide, slide_height=prs.slide_height)
    assert score.performance == 2


def test_badly_broken_slide_scores_low_across_axes():
    prs = Presentation(str(FIXTURES / "badly_broken.pptx"))
    slide = prs.slides[1]  # fragment + no annotation + no note
    score = score_slide(slide, slide_height=prs.slide_height)
    assert score.consultative <= 1
    assert score.performance <= 1


def test_score_fields_present():
    prs = Presentation(str(FIXTURES / "pristine.pptx"))
    slide = prs.slides[0]
    score = score_slide(slide, slide_height=prs.slide_height)
    assert isinstance(score, NarrativeScore)
    assert 0 <= score.consultative <= 3
    assert 0 <= score.performance <= 3
    assert 0 <= score.focal <= 3
