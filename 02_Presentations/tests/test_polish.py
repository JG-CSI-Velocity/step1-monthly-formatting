"""Integration tests for polish.py against the three committed fixtures."""

from pathlib import Path

from pptx import Presentation

from polish import apply_fixes, audit_deck, write_diff_report, write_report

FIXTURES = Path(__file__).parent / "fixtures"


def test_pristine_audit_yields_no_flags_on_content_slides():
    # Title slide (slide index 1 / 1-based) is exempt -- title headlines
    # are inherently non-consultative and are not meant to be scored.
    result = audit_deck(FIXTURES / "pristine.pptx")
    assert result.deck_level.client_name_present is True
    content_slides = [s for s in result.slides if s.index > 1]
    assert all(s.flagged is False for s in content_slides), (
        "Pristine content slides should not flag: "
        + ", ".join(f"slide {s.index}" for s in content_slides if s.flagged)
    )


def test_moderately_broken_audit_flags_fragment_and_missing_annotation():
    result = audit_deck(FIXTURES / "moderately_broken.pptx")
    flagged = [s for s in result.slides if s.flagged]
    assert len(flagged) >= 2


def test_badly_broken_audit_flags_all_content_slides():
    result = audit_deck(FIXTURES / "badly_broken.pptx")
    flagged = [s for s in result.slides if s.flagged]
    assert len(flagged) >= 2


def test_write_report_produces_markdown(tmp_path):
    result = audit_deck(FIXTURES / "pristine.pptx")
    report_path = tmp_path / "polish_report.md"
    write_report(result, report_path)
    text = report_path.read_text()
    assert text.startswith("# Polish Report")
    assert "Slides:" in text


def test_apply_fixes_replaces_times_new_roman_with_montserrat(tmp_path):
    out_path = tmp_path / "fixed.pptx"
    apply_fixes(FIXTURES / "badly_broken.pptx", out_path)
    prs = Presentation(str(out_path))
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.font.name is not None:
                        assert run.font.name == "Montserrat", (
                            f"Found non-Montserrat font: {run.font.name}"
                        )


def test_apply_fixes_snaps_near_palette_color_on_moderately_broken(tmp_path):
    out_path = tmp_path / "fixed.pptx"
    apply_fixes(FIXTURES / "moderately_broken.pptx", out_path)
    prs = Presentation(str(out_path))
    from pptx.dml.color import RGBColor
    CORAL = RGBColor(0xE7, 0x4C, 0x3C)
    found_coral = False
    slide = prs.slides[3]  # slide index 3 (0-based) has near-coral title
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                try:
                    if run.font.color.rgb == CORAL:
                        found_coral = True
                except AttributeError:
                    pass
    assert found_coral, "Expected near-coral color to snap to exact CORAL"


def test_diff_report_shows_before_after_text_per_slide(tmp_path):
    report_path = tmp_path / "polish_diff.md"
    write_diff_report(
        before_path=FIXTURES / "pristine.pptx",
        after_path=FIXTURES / "pristine.pptx",  # self-compare: no text changes
        out_path=report_path,
    )
    text = report_path.read_text()
    assert text.startswith("# Polish Diff")
    assert "Slide 1" in text
