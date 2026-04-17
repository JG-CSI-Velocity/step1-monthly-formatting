"""Generate test fixture PPTX files for the polish test suite.

Run once manually; commit the generated .pptx files:
    python 02_Presentations/scripts/make_fixtures.py

Produces three decks with known compliance properties so tests can
assert exact flag counts without depending on client data.
"""

from __future__ import annotations

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent.parent))

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from style.palette import CORAL, NAVY, TEAL

FIXTURES = Path(__file__).parent.parent / "tests" / "fixtures"


def _new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    return prs


def _add_title_slide(prs: Presentation, title: str, client: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(1.0))
    tf = tb.text_frame
    tf.text = title
    for run in tf.paragraphs[0].runs:
        run.font.name = "Montserrat"
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = NAVY
    sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.33), Inches(0.6))
    sub.text_frame.text = client
    for run in sub.text_frame.paragraphs[0].runs:
        run.font.name = "Montserrat"
        run.font.size = Pt(18)


def _add_content_slide(
    prs: Presentation,
    headline: str,
    annotation: str | None,
    speaker_note: str | None,
    font_name: str = "Montserrat",
    title_color: RGBColor = NAVY,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.4), Inches(12.33), Inches(1.0)
    )
    title_box.text_frame.text = headline
    for run in title_box.text_frame.paragraphs[0].runs:
        run.font.name = font_name
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = title_color

    if annotation is not None:
        ann_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(7.1), Inches(12.33), Inches(0.3)
        )
        ann_box.text_frame.text = annotation
        for run in ann_box.text_frame.paragraphs[0].runs:
            run.font.name = font_name
            run.font.size = Pt(9)

    if speaker_note is not None:
        slide.notes_slide.notes_text_frame.text = speaker_note


def build_pristine() -> Presentation:
    prs = _new_deck()
    _add_title_slide(prs, "Q1 2026 Debit Card Performance", "Cape & Coast Bank")
    _add_content_slide(
        prs,
        headline="Three branches drove 62% of debit growth, led by Main Office at 23%.",
        annotation="Recommend targeted staff training at top three locations.",
        speaker_note="Focus discussion on branch-level replicability.",
    )
    _add_content_slide(
        prs,
        headline="Interchange revenue rose 14% year-over-year on PIN-signature mix shift.",
        annotation="Continue promoting signature transactions through rewards.",
        speaker_note="Signature carries 2.1x the interchange of PIN.",
    )
    _add_content_slide(
        prs,
        headline="Attrition declined 8% after the Q4 re-engagement campaign.",
        annotation="Extend campaign cadence to quarterly in 2026.",
        speaker_note="Responder lift: $128 avg incremental debit spend per account.",
    )
    _add_content_slide(
        prs,
        headline="Payroll direct-deposit penetration reached 71%, up 4 points.",
        annotation="PFI status for payroll accounts is 3x portfolio average.",
        speaker_note="Upsell opportunity: auto-enroll overdraft protection.",
    )
    return prs


def build_moderately_broken() -> Presentation:
    prs = _new_deck()
    _add_title_slide(prs, "Q1 2026 Debit Card Performance", "Cape & Coast Bank")
    # Fragment headline (non-sentence)
    _add_content_slide(
        prs,
        headline="Debit Card Performance by Branch",
        annotation="Main Office leads growth.",
        speaker_note="Speaker note present.",
    )
    # Missing annotation
    _add_content_slide(
        prs,
        headline="Interchange revenue rose 14% year-over-year on PIN-signature mix shift.",
        annotation=None,
        speaker_note="PIN-signature mix comment.",
    )
    # Off-palette coral title color (near-match, should snap to CORAL)
    _add_content_slide(
        prs,
        headline="Attrition declined 8% after the Q4 re-engagement campaign.",
        annotation="Extend cadence.",
        speaker_note="Note.",
        title_color=RGBColor(0xE5, 0x4E, 0x3E),  # within threshold of CORAL
    )
    return prs


def build_badly_broken() -> Presentation:
    prs = _new_deck()
    # Title slide uses wrong font + short title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.33), Inches(1.0))
    tb.text_frame.text = "Quarterly Review"
    for run in tb.text_frame.paragraphs[0].runs:
        run.font.name = "Times New Roman"
        run.font.size = Pt(36)
    # Content slide: Times New Roman, fragment headline, no annotation, no note
    _add_content_slide(
        prs,
        headline="Debit Trends",
        annotation=None,
        speaker_note=None,
        font_name="Times New Roman",
        title_color=RGBColor(0xAB, 0x12, 0xCD),  # wild off-palette
    )
    # Another Times New Roman slide with fragment
    _add_content_slide(
        prs,
        headline="Attrition Chart",
        annotation=None,
        speaker_note="",
        font_name="Times New Roman",
    )
    return prs


def main() -> None:
    FIXTURES.mkdir(parents=True, exist_ok=True)
    build_pristine().save(FIXTURES / "pristine.pptx")
    build_moderately_broken().save(FIXTURES / "moderately_broken.pptx")
    build_badly_broken().save(FIXTURES / "badly_broken.pptx")
    print(f"Wrote 3 fixtures to {FIXTURES}")


if __name__ == "__main__":
    main()
