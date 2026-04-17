"""Consultative/performance/focal scorer for slides.

Polish flags slides with any axis < 2. No content rewriting.
"""

from __future__ import annotations

from dataclasses import dataclass

from pptx.util import Emu

from style.headline import score_headline
from style.layout import FOOTER_ZONE, TITLE_ZONE, is_inside_zone


@dataclass
class NarrativeScore:
    consultative: int  # 0-3 based on headline quality
    performance: int   # 0-3 based on annotation + speaker note presence
    focal: int         # 0-3 based on count of bold colors on the slide


def _find_title_text(slide, slide_height: Emu) -> str:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if is_inside_zone(shape, TITLE_ZONE, slide_height):
            text = shape.text_frame.text.strip()
            if text:
                return text
    return ""


def _find_footer_text(slide, slide_height: Emu) -> str:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if is_inside_zone(shape, FOOTER_ZONE, slide_height):
            text = shape.text_frame.text.strip()
            if text:
                return text
    return ""


def _speaker_note_text(slide) -> str:
    if not slide.has_notes_slide:
        return ""
    return slide.notes_slide.notes_text_frame.text.strip()


def _consultative_score(headline: str) -> int:
    if not headline:
        return 0
    s = score_headline(headline)
    passed = sum(
        [s.is_complete_sentence, s.has_metric, s.has_direction, s.has_driver_clause]
    )
    if passed == 4:
        return 3
    if passed == 3:
        return 2
    if passed >= 1:
        return 1
    return 0


def _performance_score(annotation: str, note: str) -> int:
    has_ann = bool(annotation)
    has_note = bool(note)
    if has_ann and has_note:
        return 3
    if has_ann or has_note:
        return 2
    return 0


def _focal_score(slide) -> int:
    """Count distinct bold fill colors on the slide. 1 -> 3; 2 -> 2; 3 -> 1; >=4 -> 0."""
    from pptx.dml.color import RGBColor
    seen: set[tuple[int, int, int]] = set()
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    try:
                        rgb = run.font.color.rgb
                    except AttributeError:
                        continue
                    if isinstance(rgb, RGBColor):
                        seen.add((rgb[0], rgb[1], rgb[2]))
    n = len(seen)
    if n <= 1:
        return 3
    if n == 2:
        return 2
    if n == 3:
        return 1
    return 0


def score_slide(slide, slide_height: Emu) -> NarrativeScore:
    headline = _find_title_text(slide, slide_height)
    annotation = _find_footer_text(slide, slide_height)
    note = _speaker_note_text(slide)

    return NarrativeScore(
        consultative=_consultative_score(headline),
        performance=_performance_score(annotation, note),
        focal=_focal_score(slide),
    )
