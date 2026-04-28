"""Montserrat type scale per SLIDE_MAPPING.md.

Each role is a factory returning a TextStyle. Callers apply the style via
`apply(text_frame, style)` which sets every run in the frame.
"""

from __future__ import annotations

from dataclasses import dataclass

from pptx.util import Pt


@dataclass(frozen=True)
class TextStyle:
    font_name: str
    size_pt: int
    weight: str  # "Regular" | "Bold" | "ExtraBold"


def slide_title() -> TextStyle:
    return TextStyle("Montserrat", 20, "ExtraBold")


def subtitle() -> TextStyle:
    return TextStyle("Montserrat", 14, "Regular")


def chart_title() -> TextStyle:
    return TextStyle("Montserrat", 18, "Bold")


def chart_subtitle() -> TextStyle:
    return TextStyle("Montserrat", 12, "Regular")


def data_label() -> TextStyle:
    return TextStyle("Montserrat", 11, "Regular")


def axis_label() -> TextStyle:
    return TextStyle("Montserrat", 10, "Regular")


def footnote() -> TextStyle:
    return TextStyle("Montserrat", 9, "Regular")


def kpi_hero() -> TextStyle:
    return TextStyle("Montserrat", 36, "ExtraBold")


def kpi_label() -> TextStyle:
    return TextStyle("Montserrat", 11, "Regular")


def apply(text_frame, style: TextStyle) -> None:
    """Set every run in `text_frame` to match `style`.

    Bold flag is True for weights "Bold" or "ExtraBold", else False.
    python-pptx does not expose an ExtraBold setter directly; the bold
    attribute paired with font_name='Montserrat' renders correctly when
    the Montserrat ExtraBold face is installed on the render host.
    """
    bold = style.weight in ("Bold", "ExtraBold")
    size = Pt(style.size_pt)
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.name = style.font_name
            run.font.size = size
            run.font.bold = bold
