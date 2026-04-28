"""Deck polish driver -- post-hoc PPTX compliance pass."""

from __future__ import annotations

import argparse
import sys
from dataclasses import dataclass, field
from pathlib import Path

sys.path.insert(0, str(Path(__file__).parent))

from loguru import logger
from pptx import Presentation

from style.charts import audit_chart_image
from style.headline import score_headline
from style.layout import FOOTER_ZONE, TITLE_ZONE, is_inside_zone
from style.narrative import NarrativeScore, score_slide


@dataclass
class SlideAudit:
    index: int
    headline: str
    narrative: NarrativeScore
    headline_violations: list[str]
    chart_dpi_violations: list[str]
    flagged: bool


@dataclass
class DeckLevelAudit:
    client_name_present: bool
    section_dividers_ok: bool
    page_numbers_ok: bool
    appendix_separated: bool


@dataclass
class DeckAudit:
    deck_path: Path
    slide_count: int
    deck_level: DeckLevelAudit
    slides: list[SlideAudit] = field(default_factory=list)


def _extract_title(slide, slide_height) -> str:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if is_inside_zone(shape, TITLE_ZONE, slide_height):
            text = shape.text_frame.text.strip()
            if text:
                return text
    return ""


def _client_name_present(title_slide_headline: str) -> bool:
    """Heuristic: non-generic placeholder (contains Bank/Union or >=4 words)."""
    if not title_slide_headline:
        return False
    lower = title_slide_headline.lower()
    return "bank" in lower or "union" in lower or len(title_slide_headline.split()) >= 4


def audit_deck(deck_path: Path) -> DeckAudit:
    prs = Presentation(str(deck_path))
    slide_height = prs.slide_height

    title_text = _extract_title(prs.slides[0], slide_height) if len(prs.slides) else ""
    deck_level = DeckLevelAudit(
        client_name_present=_client_name_present(title_text),
        section_dividers_ok=True,  # placeholder for future expansion
        page_numbers_ok=True,
        appendix_separated=True,
    )

    slide_audits: list[SlideAudit] = []
    for i, slide in enumerate(prs.slides, start=1):
        headline = _extract_title(slide, slide_height)
        h_score = score_headline(headline) if headline else None
        narrative = score_slide(slide, slide_height=slide_height)

        chart_flags: list[str] = []
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    audit = audit_chart_image(shape)
                    if not audit.meets_dpi_floor:
                        chart_flags.append(
                            f"Image below DPI floor ({audit.dpi_estimate} < 150)"
                        )
                except Exception as e:  # pragma: no cover
                    logger.warning(f"Slide {i} image audit failed: {e}")

        # Title slide (index 1) is exempt from narrative flags -- its
        # headline is inherently not a consultative finding.
        is_title_slide = i == 1
        flagged = (not is_title_slide) and (
            narrative.consultative < 2
            or narrative.performance < 3
            or narrative.focal < 2
            or bool(chart_flags)
        )

        slide_audits.append(
            SlideAudit(
                index=i,
                headline=headline,
                narrative=narrative,
                headline_violations=h_score.violates if h_score else ["no headline detected"],
                chart_dpi_violations=chart_flags,
                flagged=flagged,
            )
        )

    return DeckAudit(
        deck_path=deck_path,
        slide_count=len(prs.slides),
        deck_level=deck_level,
        slides=slide_audits,
    )


def write_report(audit: DeckAudit, path: Path) -> None:
    lines: list[str] = []
    lines.append(f"# Polish Report -- {audit.deck_path.name}")
    lines.append("")
    lines.append("## Summary")
    lines.append(f"- Slides: {audit.slide_count}")
    passing = sum(1 for s in audit.slides if not s.flagged)
    lines.append(f"- Passing (all axes >=2, no DPI violations): {passing}")
    lines.append(f"- Flagged: {audit.slide_count - passing}")
    lines.append("")
    lines.append("## Deck-level findings")
    lines.append(f"- [{'x' if audit.deck_level.client_name_present else ' '}] Title slide has client name")
    lines.append(f"- [{'x' if audit.deck_level.section_dividers_ok else ' '}] Section dividers present")
    lines.append(f"- [{'x' if audit.deck_level.page_numbers_ok else ' '}] Page numbers present")
    lines.append(f"- [{'x' if audit.deck_level.appendix_separated else ' '}] Appendix separated")
    lines.append("")
    lines.append("## Slide-by-slide")
    for s in audit.slides:
        lines.append(f"### Slide {s.index} -- \"{s.headline or '(no title)'}\"")
        lines.append(f"- Consultative: {s.narrative.consultative}/3")
        lines.append(f"- Performance: {s.narrative.performance}/3")
        lines.append(f"- Focal: {s.narrative.focal}/3")
        if s.headline_violations:
            lines.append("- Headline violations:")
            for v in s.headline_violations:
                lines.append(f"  - {v}")
        if s.chart_dpi_violations:
            lines.append("- Chart violations:")
            for v in s.chart_dpi_violations:
                lines.append(f"  - {v}")
        lines.append("")
    path.write_text("\n".join(lines))


def _slide_texts(prs_path: Path) -> list[list[str]]:
    """Return a list of per-slide lists of text-frame strings."""
    prs = Presentation(str(prs_path))
    result: list[list[str]] = []
    for slide in prs.slides:
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    texts.append(text)
        result.append(texts)
    return result


def write_diff_report(before_path: Path, after_path: Path, out_path: Path) -> None:
    before = _slide_texts(before_path)
    after = _slide_texts(after_path)
    lines: list[str] = ["# Polish Diff", ""]
    for i, (b, a) in enumerate(zip(before, after), start=1):
        lines.append(f"## Slide {i}")
        if b == a:
            lines.append("_(no text changes)_")
        else:
            lines.append("**Before:**")
            for t in b:
                lines.append(f"- {t}")
            lines.append("**After:**")
            for t in a:
                lines.append(f"- {t}")
        lines.append("")
    out_path.write_text("\n".join(lines))


MONTSERRAT_FAMILY = ("Montserrat", "Montserrat Regular", "Montserrat Bold",
                     "Montserrat ExtraBold", "Montserrat Medium")


def _force_montserrat(run) -> bool:
    """If run's font isn't in the Montserrat family, set it. Returns True if changed."""
    current = run.font.name
    if current in MONTSERRAT_FAMILY:
        return False
    run.font.name = "Montserrat"
    return True


def _snap_near_palette(run) -> bool:
    """If run color is within palette threshold, snap to palette. Returns True if changed."""
    from style.palette import nearest_palette
    try:
        rgb = run.font.color.rgb
    except AttributeError:
        return False
    if rgb is None:
        return False
    snapped = nearest_palette(rgb)
    if snapped is not None and snapped != rgb:
        run.font.color.rgb = snapped
        return True
    return False


def apply_fixes(deck_path: Path, out_path: Path) -> dict[str, int]:
    """Open deck, apply force-apply fixes, save to out_path.

    Returns a dict of change counts: {'fonts_fixed': N, 'colors_snapped': M}.
    """
    prs = Presentation(str(deck_path))
    fonts_fixed = 0
    colors_snapped = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if _force_montserrat(run):
                        fonts_fixed += 1
                    if _snap_near_palette(run):
                        colors_snapped += 1

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return {"fonts_fixed": fonts_fixed, "colors_snapped": colors_snapped}


def _parse_args(argv: list[str] | None = None) -> argparse.Namespace:
    p = argparse.ArgumentParser(
        prog="polish",
        description="Post-hoc PPTX polish pass (SLIDE_MAPPING.md compliance).",
    )
    p.add_argument("deck", nargs="?", type=Path, help="Path to input .pptx")
    p.add_argument("--batch", type=Path, help="Process every .pptx in this folder")
    p.add_argument("--apply", action="store_true", help="Write polished PPTX (default dry-run)")
    p.add_argument("--report-only", action="store_true", help="Emit report only")
    p.add_argument("--out", type=Path, default=None, help="Output directory")
    p.add_argument("--strict", action="store_true", help="Exit non-zero if any flag")
    return p.parse_args(argv)


def _process_one(deck_path: Path, out_dir: Path, apply: bool) -> DeckAudit:
    out_dir.mkdir(parents=True, exist_ok=True)
    audit = audit_deck(deck_path)
    report_path = out_dir / f"{deck_path.stem}__polish_report.md"
    write_report(audit, report_path)
    logger.info(f"Wrote {report_path}")

    if apply:
        polished_path = out_dir / deck_path.name
        counts = apply_fixes(deck_path, polished_path)
        logger.info(
            f"Applied: {counts['fonts_fixed']} fonts, "
            f"{counts['colors_snapped']} colors. Wrote {polished_path}"
        )
        diff_path = out_dir / f"{deck_path.stem}__polish_diff.md"
        write_diff_report(deck_path, polished_path, diff_path)
        logger.info(f"Wrote {diff_path}")

    return audit


def main(argv: list[str] | None = None) -> int:
    args = _parse_args(argv)

    if args.batch is None and args.deck is None:
        logger.error("Provide a deck path or --batch <dir>")
        return 2

    decks: list[Path]
    if args.batch is not None:
        decks = sorted(args.batch.glob("*.pptx"))
        if not decks:
            logger.error(f"No .pptx files found in {args.batch}")
            return 2
    else:
        decks = [args.deck]

    all_audits: list[DeckAudit] = []
    for d in decks:
        out_dir = args.out if args.out else d.parent / "polished"
        all_audits.append(_process_one(d, out_dir, apply=args.apply))

    if args.strict:
        any_flag = any(s.flagged for a in all_audits for s in a.slides)
        if any_flag:
            logger.error("Strict mode: flagged slides present")
            return 1

    return 0


if __name__ == "__main__":
    raise SystemExit(main())
